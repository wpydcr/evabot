import os
from pathlib import Path
import platform
import re
import shutil
import subprocess
import tempfile
from typing import Any,List, Dict
import requests
from markdownify import markdownify as md
import os
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import pandas as pd
import locale
system_encoding = locale.getpreferredencoding()

def _guard_command(command: str, cwd: str) -> str | None:
    cmd = command.strip()
    lower = cmd.lower()
    deny_patterns = [
            r"\brm\s+-[rf]{1,2}\b",          # rm -r, rm -rf, rm -fr
            r"\bdel\s+/[fq]\b",              # del /f, del /q
            r"\brmdir\s+/s\b",               # rmdir /s
            r"\b(format|mkfs|diskpart)\b",   # disk operations
            r"\bdd\s+if=",                   # dd
            r">\s*/dev/sd",                  # write to disk
            r"\b(shutdown|reboot|poweroff)\b",  # system power
            r":\(\)\s*\{.*\};\s*:",          # fork bomb
        ]
    for pattern in deny_patterns:
        if re.search(pattern, lower):
            return "Error: Command blocked by safety guard (dangerous command detected)"

    return None

def _resolve_path(path: str, allowed_dir: Path | None = None) -> Path:
    """Resolve path and optionally enforce directory restriction."""
    resolved = Path(path).expanduser().resolve()
    # if allowed_dir and not str(resolved).startswith(str(allowed_dir.resolve())):
    #     raise PermissionError(f"Path {path} is outside allowed directory {allowed_dir}")
    return resolved

def read_pdf(file_path: str) -> str:
    text = ""
    try:
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text.strip()
    except Exception as e:
        return f"[读取PDF失败]: {str(e)}"

def read_docx(file_path: str) -> str:
    try:
        doc = docx.Document(file_path)
        full_text = [para.text for para in doc.paragraphs if para.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        full_text.append(cell.text.strip())
        return "\n".join(full_text)
    except Exception as e:
        return f"[读取DOCX失败]: {str(e)}"

def read_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text.strip())
        return "\n".join(text_runs)
    except Exception as e:
        return f"[读取PPTX失败]: {str(e)}"

def read_excel(file_path: str) -> str:
    """读取 Excel 文件 (xlsx/xls)，提取所有 Sheet 的数据"""
    try:
        xls = pd.ExcelFile(file_path)
        text_output = []
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            if not df.empty:
                text_output.append(f"\n--- Sheet: {sheet_name} ---")
                # 转化为CSV格式的字符串，对LLM非常友好
                text_output.append(df.to_csv(index=False))
        return "\n".join(text_output).strip()
    except Exception as e:
        return f"[读取Excel失败]: {str(e)}"

def get_libreoffice_command() -> str:
    """获取当前系统下的 LibreOffice 调用命令"""
    system = platform.system()
    if system == "Windows":
        return "soffice" # 确保 LibreOffice 安装目录已加入 PATH
    elif system == "Darwin": # macOS
        return "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    else: # Linux
        return "libreoffice"

def read_legacy_office(file_path: str, ext: str) -> str:
    """
    通过 LibreOffice 将 .doc/.ppt 转换为 .docx/.pptx 临时文件，然后读取
    """
    target_ext = "docx" if ext == ".doc" else "pptx"
    cmd = get_libreoffice_command()
    
    # 创建一个临时目录存放转换后的文件，避免污染原目录
    temp_dir = tempfile.mkdtemp()
    
    try:
        # 组装转换命令：libreoffice --headless --convert-to docx/pptx --outdir <temp_dir> <file_path>
        process = subprocess.run(
            [cmd, "--headless", "--convert-to", target_ext, "--outdir", temp_dir, file_path],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True
        )
        
        if process.returncode != 0:
            return f"[格式转换失败]: 请检查是否安装了 LibreOffice。详细错误: {process.stderr}"
        
        # 获取转换后的新文件路径
        base_name = os.path.basename(file_path)
        new_file_name = os.path.splitext(base_name)[0] + f".{target_ext}"
        new_file_path = os.path.join(temp_dir, new_file_name)
        
        if not os.path.exists(new_file_path):
            return "[格式转换失败]: 未能生成目标文件。"
            
        # 调用现有的现代格式读取函数
        if target_ext == "docx":
            return read_docx(new_file_path)
        else:
            return read_pptx(new_file_path)
            
    except FileNotFoundError:
        return "[系统错误]: 找不到 LibreOffice 命令行工具，请确保已安装并配置环境变量。"
    except Exception as e:
        return f"[读取旧版Office文件失败]: {str(e)}"
    finally:
        # 清理临时目录和转换产生的文件
        shutil.rmtree(temp_dir, ignore_errors=True)

def read_local_file(file_path: str) -> str:
    """统一的文件读取入口"""
    if not os.path.exists(file_path):
        return f"[错误]: 文件不存在 '{file_path}'"

    ext = os.path.splitext(file_path)[1].lower()

    if ext == '.pdf':
        return read_pdf(file_path)
    elif ext in ['.docx']:
        return read_docx(file_path)
    elif ext in ['.pptx']:
        return read_pptx(file_path)
    elif ext in ['.xlsx', '.xls']:
        return read_excel(file_path)
    elif ext in ['.doc', '.ppt']:
        return read_legacy_office(file_path, ext)
    else:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            return f"[读取文件失败]: {str(e)}"

# ==========================================
# 2. 面向对象的工具类封装
# ==========================================

class BaseTool:
    """所有工具的基类"""
    name: str = ""
    description: str = ""
    parameters: Dict[str, Any] = {}

    def execute(self, **kwargs: Any) -> str:
        raise NotImplementedError("Subclasses must implement the execute method.")


class ListDirTool(BaseTool):
    name = "list_dir"
    description = "List the contents of a directory."
    parameters = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "The directory path to list"}
        },
        "required": ["path"]
    }

    def execute(self, path: str, **kwargs: Any) -> str:
        try:
            dir_path = _resolve_path(path)
            if not dir_path.exists():
                return f"Error: Directory not found: {path}"
            if not dir_path.is_dir():
                return f"Error: Not a directory: {path}"
            
            items = []
            for item in sorted(dir_path.iterdir()):
                prefix = "[DIR] " if item.is_dir() else "[FILE]"
                if prefix == "[FILE]":
                    size_mb = item.stat().st_size / (1024 * 1024)
                    items.append(f"{prefix}{item.name} ({size_mb:.1f} MB)")
                else:
                    items.append(f"{prefix}{item.name}")
            
            if not items:
                return f"Directory {path} is empty"
            return "\n".join(items)
        except PermissionError as e:
            return f"Error: {e}"
        except Exception as e:
            return f"Error listing directory: {str(e)}"


class ExecCommandTool(BaseTool):
    name = "exec_command"
    description = "Execute a shell command and return its output. Use with caution."
    parameters = {
        "type": "object",
        "properties": {
            "command": {"type": "string", "description": "The shell command to execute"},
            "working_dir": {"type": "string", "description": "Optional working directory for the command"}
        },
        "required": ["command"]
    }
    
    def execute(self, command: str, working_dir: str = None, **kwargs: Any) -> str:
        cwd = working_dir if working_dir else os.getcwd()
        guard_error = _guard_command(command, cwd)
        if guard_error:
            return guard_error
        
        try:
            process = subprocess.run(
                command,
                shell=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                cwd=cwd,
                timeout=300  
            )
            
            output_parts = []
            if process.stdout:
                try:
                    # 先尝试统一的 utf-8 解码
                    stdout_text = process.stdout.decode("utf-8")
                except UnicodeDecodeError:
                    # 如果报错，尝试使用系统默认编码 (如 Windows 的 GBK) 解码
                    stdout_text = process.stdout.decode(system_encoding, errors="replace")
                output_parts.append(stdout_text)
            if process.stderr:
                try:
                    stderr_text = process.stderr.decode("utf-8")
                except UnicodeDecodeError:
                    stderr_text = process.stderr.decode(system_encoding, errors="replace")
                if stderr_text.strip():
                    output_parts.append(f"STDERR:\n{stderr_text.strip()}")
            if process.returncode != 0:
                output_parts.append(f"\nExit code: {process.returncode}")
            
            result = "\n".join(output_parts) if output_parts else "(no output)"
            max_len = 50000
            if len(result) > max_len:
                result = result[:max_len] + f"\n... (truncated, {len(result) - max_len} more chars)"
            return result
            
        except subprocess.TimeoutExpired:
            return "Error: Command timed out after 60 seconds"
        except Exception as e:
            return f"Error executing command: {str(e)}"

class ReadFileTool(BaseTool):
    name = "read_file"
    description = "Read the contents of a file at the given path."
    parameters = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "The file path to read"}
        },
        "required": ["path"]
    }

    def execute(self, path: str, **kwargs: Any) -> str:
        try:
            file_path = _resolve_path(path)
            if not file_path.exists():
                return f"Error: File not found: {path}"
            if not file_path.is_file():
                return f"Error: Not a file: {path}"
            return read_local_file(str(file_path))
        except PermissionError as e:
            return f"Error: {e}"
        except Exception as e:
            return f"Error reading file: {str(e)}"

class WriteFileTool(BaseTool):
    name = "write_file"
    description = "Write content to a plain text file at the given path. Creates parent directories if needed."
    parameters = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "The file path to write to"},
            "content": {"type": "string", "description": "The content to write"}
        },
        "required": ["path", "content"]
    }

    def execute(self, path: str, content: str, **kwargs: Any) -> str:
        try:
            file_path = _resolve_path(path)
            file_path.parent.mkdir(parents=True, exist_ok=True)
            file_path.write_text(content, encoding="utf-8")
            return f"Successfully wrote {len(content)} bytes to {path}"
        except PermissionError as e:
            return f"Error: {e}"
        except Exception as e:
            return f"Error writing file: {str(e)}"

class EditFileTool(BaseTool):
    name = "edit_file"
    description = "Edit a plain text file by replacing old_text with new_text. The old_text must exist exactly in the file."
    parameters = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "The file path to edit"},
            "old_text": {"type": "string", "description": "The exact text to find and replace"},
            "new_text": {"type": "string", "description": "The text to replace with"}
        },
        "required": ["path", "old_text", "new_text"]
    }

    def execute(self, path: str, old_text: str, new_text: str, **kwargs: Any) -> str:
        try:
            file_path = _resolve_path(path)
            if not file_path.exists():
                return f"Error: File not found: {path}"
            
            content = file_path.read_text(encoding="utf-8")
            
            if old_text not in content:
                return "Error: old_text not found in file. Make sure it matches exactly."
            
            count = content.count(old_text)
            if count > 1:
                return f"Warning: old_text appears {count} times. Please provide more context to make it unique."
            
            new_content = content.replace(old_text, new_text, 1)
            file_path.write_text(new_content, encoding="utf-8")
            
            return f"Successfully edited {path}"
        except PermissionError as e:
            return f"Error: {e}"
        except Exception as e:
            return f"Error editing file: {str(e)}"

class GetAllSkillsTool(BaseTool):
    name = "get_all_skills"
    description = "获取系统中所有 skills 的树状结构，包含技能名称、所在目录路径和功能描述。"
    parameters = {
        "type": "object",
        "properties": {},
        "required": []
    }

    def execute(self, **kwargs: Any) -> str:
        try:
            # 局部导入，避免顶层导入可能导致的循环引用
            from backend.power.power import PowerManager
            pm = PowerManager()
            
            def build_tree(skills_dict, indent=0):
                lines = []
                prefix = "  " * indent
                for name, skill in skills_dict.items():
                    # 剔除换行符，保证一行一条，最省 token
                    desc = skill.description.replace('\n', ' ').strip()
                        
                    # 极简树状表达格式: - name: description
                    lines.append(f"{prefix}- {skill.name}: {desc}")
                    
                    # 递归追加子技能
                    if skill.sub_skills:
                        lines.extend(build_tree(skill.sub_skills, indent + 1))
                return lines
            
            tree_lines = build_tree(pm.skills)
            if not tree_lines:
                return "暂无可用的 skills。"
            return "\n".join(tree_lines)
        except Exception as e:
            return f"获取 skills 失败: {str(e)}"

# 注册所有支持的工具
REGISTERED_TOOLS = [
    ListDirTool(),
    ExecCommandTool(),
    ReadFileTool(),
    WriteFileTool(),
    EditFileTool(),
    GetAllSkillsTool()
]


task_tool_schema = [{
    "type": "function", 
    "function": {
        "name": 'use_skill',
        "description": 'Use the appropriate skill to complete a task.',
        "parameters": {
        "type": "object",
        "properties": {
            "skill_name": {"type": "string", "description": "The name of the skill to use"},
            "goal": {"type": "string", "description": "The specific goal for using this skill, providing clear instructions including needed information and expected outcomes"},
            "needs_self_verification": {
                "type": "boolean",
                "description": "True if the task involves operations like file-changing (e.g., creating/modifying files, open APP) that require explicit verification of the outcome. False for read-only or informational queries (e.g., searching web, reading file content)."
            }
        },
        "required": ["skill_name","goal","needs_self_verification"]
    }
    }
},
{
    "type": "function", 
    "function": {
        "name": "communicate_with_upstream",
        "description": "Communicate with the task source (upstream). Use this tool to request missing information needed to clarify the task intent, OR to answer questions when queried by the upstream source.",
        "parameters": {
            "type": "object",
            "properties": {
                "send_info": {
                    "type": "string", 
                    "description": "The message to send to the task source. This can be a request for specific missing details, or your answer to their inquiry."
                }
            },
            "required": ["send_info"]
        }
    }
},
{
    "type": "function",
    "function": {
        "name": "communicate_with_downstream",
        "description": "Communicate with a downstream task executor. Use this tool to provide important new information, answer the executor's questions, ask follow-up questions, or assign subsequent tasks to re-awaken the executor.",
        "parameters": {
            "type": "object",
            "properties": {
                "tool_call_id": {
                    "type": "string", 
                    "description": "The downstream task executor's tool_call_id."
                },
                "provide_info": {
                    "type": "string", 
                    "description": "The message to send to the executor. This can be an answer to their question, new context, follow-up questions, or instructions for a new sub-task. DO NOT put the executor's original question here."
                }
            },
            "required": ["tool_call_id", "provide_info"]
        }
    }
},
{
    "type": "function",
    "function": {
        "name": "report_deliverable_file",
        "description": "Report a deliverable file to the upstream and mark the current task as completed. Call this tool ONLY when the task requires delivering a specific file. If no file needs to be delivered, do not call this tool; instead, output your final response directly to complete the task.",
        "parameters": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "The local path of the deliverable file, which should be an absolute path."},
                "description": {"type": "string", "description": "A brief description of the file"}
            },
            "required": ["file_path", "description"]
        }
    }
}
]

# ==========================================
# 3. 统一入口函数
# ==========================================

def get_base_tool(tool_name=[]) -> List[Dict[str, Any]]:
    base_tool_schema = []
    for tool in REGISTERED_TOOLS:
        if tool.name in tool_name or tool_name == []:
            base_tool_schema.append({
                "type": "function",
                "function": {
                    "name": tool.name,
                    "description": tool.description,
                    "parameters": tool.parameters
                }
            })
    for tool in task_tool_schema:
        if tool["function"]["name"] in tool_name or tool_name == []:
            base_tool_schema.append(tool)
    return base_tool_schema

def execute_tool(name: str, params: dict[str, Any]) -> str:
    """路由并执行特定的工具"""
    for tool in REGISTERED_TOOLS:
        if tool.name == name:
            try:
                return tool.execute(**params)
            except Exception as e:
                return f"Error executing {name}: {str(e)}"
    return f"Error: Tool '{name}' not found"



